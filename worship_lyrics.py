from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from warnings import warn
import os.path as path
from numbers import Real as RealNumbers
from packaging import version

__version__ = version.parse("2.1.0")


class ParseError(RuntimeError):
    def __init__(self, line, message):
        real_msg = "On line {}: {}".format(line, message)
        super().__init__(real_msg)


class ParseWarning(Warning):
    def __init__(self, message):
        real_msg = "On line {}: {}".format(message[0], message[1])
        super().__init__(real_msg)


def initialize_pptx(width: RealNumbers = 16, height: RealNumbers = 9) -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)

    return prs


def create_slide(prs: Presentation, curr_line, curr_lyrics, vars):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    # background image
    if "BKG" in vars and vars["BKG"] != "":
        shapes.add_picture(
            vars["BKG"], 0, 0,
            height=prs.slide_height, width=prs.slide_width
        )
    else:
        warn((curr_line, "No background image while creating slides"), ParseWarning)
    # lyrics text
    if curr_lyrics != "":
        text_shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(vars["SLIDE-HEIGHT"] / 6),
            Inches(vars["SLIDE-WIDTH"]),
            Inches(vars["SLIDE-HEIGHT"] - (vars["SLIDE-HEIGHT"] / 6))
        )
        text_shape.fill.background()
        text_shape.line.fill.background()
        text_frame = text_shape.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_paragraph = text_frame.paragraphs[0]
        text_paragraph.text = curr_lyrics
        text_paragraph.alignment = PP_ALIGN.CENTER
        # lyrics font size
        if "LYRICS-SIZE" in vars and vars["LYRICS-SIZE"] != "":
            try:
                font_size = float(vars["LYRICS-SIZE"])
            except ValueError as e:
                raise ParseError(curr_line,
                                 "Invalid lyrics size {}".format(vars["LYRICS-SIZE"])) from e
            text_paragraph.font.size = Pt(font_size)
        else:
            warn((curr_line, "No lyrics font size while creating slides"), ParseWarning)
        # lyrics font name
        if "LYRICS-FONT" in vars and vars["LYRICS-FONT"] != "":
            font_name = vars["LYRICS-FONT"]
            text_paragraph.font.name = font_name
        else:
            warn((curr_line, "No lyrics font name while creating slides"), ParseWarning)
        # lyrics font color
        if "LYRICS-COLOR" in vars and vars["LYRICS-COLOR"] != "":
            try:
                font_color = RGBColor.from_string(vars["LYRICS-COLOR"])
            except ValueError as e:
                raise ParseError(curr_line,
                                 "Invalid lyrics color {}".format(vars["LYRICS-COLOR"])) from e
            text_paragraph.font.color.rgb = font_color
        text_paragraph.font.bold = True
    # footer
    if "FOOTER" in vars and vars["FOOTER"] != "":
        text_shape = shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(vars["SLIDE-WIDTH"] / 32),
            Inches(vars["SLIDE-HEIGHT"] - (vars["SLIDE-WIDTH"] / 32) - 2),
            Inches(vars["SLIDE-WIDTH"] - (vars["SLIDE-WIDTH"] / 32)), Inches(2)
        )
        text_shape.fill.background()
        text_shape.line.fill.background()
        text_frame = text_shape.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_ANCHOR.BOTTOM
        text_paragraph = text_frame.paragraphs[0]
        text_paragraph.text = vars["FOOTER"]
        text_paragraph.alignment = PP_ALIGN.LEFT
        # footer font size
        if "FOOTER-SIZE" in vars and vars["FOOTER-SIZE"] != "":
            try:
                font_size = float(vars["FOOTER-SIZE"])
            except ValueError as e:
                raise ParseError(curr_line,
                                 "Invalid footer size {}".format(vars["FOOTER-SIZE"])) from e
            text_paragraph.font.size = Pt(font_size)
        else:
            warn((curr_line, "No footer font size while creating slides"), ParseWarning)
        # footer font name
        if "FOOTER-FONT" in vars and vars["FOOTER-FONT"] != "":
            font_name = vars["FOOTER-FONT"]
            text_paragraph.font.name = font_name
        else:
            warn((curr_line, "No footer font name while creating slides"), ParseWarning)
        # footer font color
        if "FOOTER-COLOR" in vars and vars["FOOTER-COLOR"] != "":
            try:
                font_color = RGBColor.from_string(vars["FOOTER-COLOR"])
            except ValueError as e:
                raise ParseError(curr_line,
                                 "Invalid footer color {}".format(vars["FOOTER-COLOR"])) from e
            text_paragraph.font.color.rgb = font_color
        text_paragraph.font.bold = True


def ensure_has_param(curr_line, cmdline, curr_state):
    if cmdline[1] == "":
        raise ParseError(curr_line,
                         "Missing argument for command {} in state {}".format(cmdline[0], curr_state))


def interpepter(filename, savefilename):
    vars = {
        "SLIDE-WIDTH":  16,
        "SLIDE-HEIGHT": 9
    }
    curr_lyrics = ""
    curr_state = "read"
    curr_line = 1
    prs = None
    file_dir = path.dirname(filename)
    with open(filename, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if curr_line == 1:
                if line[0:5] != "!VER ":
                    raise ParseError(1,
                                     "First line of lyrics file must be \"!VER\" command, got \"{}\" instead".format(line))
                version_string = line[5:]
                try:
                    ver = version.parse(version_string)
                except version.InvalidVersion as e:
                    raise ParseError(1,
                                     "Invalid version \"{}\"".format(version_string)) from e
                if ver.major != __version__.major:
                    raise ParseError(1,
                                     "Major version \"{}\" in the lyrics file does not match that of this script ({})".format(ver.major, __version__.major))
                if ver.minor > __version__.minor:
                    raise ParseError(1,
                                     "Minor version \"{}\" in the lyrics file is newer than that of this script ({})".format(ver.minor, __version__.minor))
                if ver.micro != __version__.micro:
                    warn((1, "Micro version \"{}\" in the lyrics file does not match that of this script ({})".format(
                        ver.micro, __version__.micro)), ParseWarning)
            elif len(line) > 0 and line[0] == "#":
                continue
            else:
                match curr_state:
                    case "read":
                        if line == "":
                            if curr_lyrics != "":
                                if prs == None:
                                    prs = initialize_pptx(
                                        vars["SLIDE-WIDTH"], vars["SLIDE-HEIGHT"])
                                create_slide(prs, curr_line, curr_lyrics, vars)
                                curr_lyrics = ""
                        elif line[0] == "!":
                            cmdline = line[1:].split(" ", 1)
                            if len(cmdline) == 1:
                                cmdline.append("")
                            cmdline[1] = cmdline[1].strip()
                            match cmdline[0]:
                                case "WIDTH":
                                    if prs != None:
                                        raise ParseError(curr_line,
                                                         "Attempt to use !WIDTH command after the first slide")
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    try:
                                        vars["SLIDE-WIDTH"] = float(cmdline[1])
                                    except ValueError as e:
                                        raise ParseError(curr_line,
                                                         "Failed to parse {} as number".format(cmdline[1])) from e
                                case "HEIGHT":
                                    if prs != None:
                                        raise ParseError(curr_line,
                                                         "Attempt to use !HEIGHT command after the first slide")
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    try:
                                        vars["SLIDE-HEIGHT"] = float(
                                            cmdline[1])
                                    except ValueError as e:
                                        raise ParseError(curr_line,
                                                         "Failed to parse {} as number".format(cmdline[1])) from e
                                case "BKG":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["BKG"] = path.join(
                                        file_dir, cmdline[1])
                                case "FONT":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-FONT"] = cmdline[1]
                                    vars["FOOTER-FONT"] = cmdline[1]
                                case "LYRICS-FONT":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-FONT"] = cmdline[1]
                                case "FOOTER-FONT":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["FOOTER-FONT"] = cmdline[1]
                                case "COLOR":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-COLOR"] = cmdline[1]
                                    vars["FOOTER-COLOR"] = cmdline[1]
                                case "LYRICS-COLOR":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-COLOR"] = cmdline[1]
                                case "FOOTER-COLOR":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["FOOTER-COLOR"] = cmdline[1]
                                case "SIZE":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-SIZE"] = cmdline[1]
                                    vars["FOOTER-SIZE"] = cmdline[1]
                                case "LYRICS-SIZE":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["LYRICS-SIZE"] = cmdline[1]
                                case "FOOTER-SIZE":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    vars["FOOTER-SIZE"] = cmdline[1]
                                case "EMPTY":
                                    if curr_lyrics != "":
                                        raise ParseError(curr_line,
                                                         "Attempt to execute !EMPTY with lyrics")
                                    if prs == None:
                                        prs = initialize_pptx(
                                            vars["SLIDE-WIDTH"], vars["SLIDE-HEIGHT"])
                                    create_slide(prs, curr_line, "", vars)
                                case "SECTION":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    if curr_lyrics != "":
                                        raise ParseError(curr_line,
                                                         "Attempt to execute !SECTION with lyrics")
                                    if ("SECTION-" + cmdline[1]) not in vars:
                                        raise ParseError(curr_line,
                                                         "Attempt to refer to invalid section {}".format(cmdline[0]))
                                    section_lyrics = tuple(
                                        s.strip() for s in vars["SECTION-" + cmdline[1]].split("\n\n"))
                                    if prs == None:
                                        prs = initialize_pptx(
                                            vars["SLIDE-WIDTH"], vars["SLIDE-HEIGHT"])
                                    for s in section_lyrics:
                                        create_slide(
                                            prs, curr_line, s, vars)
                                case "FOOTER-START":
                                    if curr_lyrics != "":
                                        raise ParseError(curr_line,
                                                         "Attempt to execute !FOOTER-START with lyrics")
                                    curr_state = "cmd-FOOTER"
                                case "SECTION-START":
                                    ensure_has_param(
                                        curr_line, cmdline, curr_state)
                                    if curr_lyrics != "":
                                        raise ParseError(curr_line,
                                                         "Attempt to execute !SECTION-START with lyrics")
                                    if ("SECTION-" + cmdline[1]) in vars:
                                        raise ParseError(curr_line,
                                                         "Attempt to redefine section {}".format(cmdline[0]))
                                    vars["CURR-SECTION"] = cmdline[1]
                                    curr_state = "cmd-SECTION"
                                case _:
                                    raise ParseError(curr_line,
                                                     "Invalid command {} in state {}".format(cmdline[0], curr_state))
                        else:
                            curr_lyrics += line + "\n"
                    case "cmd-FOOTER":
                        if len(line) > 0 and line[0] == "!":
                            if line == "!FOOTER-END":
                                vars["FOOTER"] = curr_lyrics
                                curr_lyrics = ""
                                curr_state = "read"
                            else:
                                raise ParseError(
                                    curr_line, "Attempt to invoke command in state {}".format(curr_state))
                        else:
                            curr_lyrics += line + "\n"
                    case "cmd-SECTION":
                        if len(line) > 0 and line[0] == "!":
                            if line == "!SECTION-END":
                                vars["SECTION-" + vars["CURR-SECTION"]
                                     ] = curr_lyrics
                                del vars["CURR-SECTION"]
                                curr_lyrics = ""
                                curr_state = "read"
                            else:
                                raise ParseError(
                                    curr_line, "Attempt to invoke command in state {}".format(curr_state))
                        else:
                            curr_lyrics += line + "\n"
            curr_line += 1
        # end for line in f
    prs.save(savefilename)


if __name__ == "__main__":
    print("Worship Lyrics version {}".format(__version__))
    import argparse
    parser = argparse.ArgumentParser(description='Create worship lyrics PPTX')
    parser.add_argument('filename')
    parser.add_argument('savefilename')
    args = parser.parse_args()

    interpepter(args.filename, args.savefilename)

    from sys import exit
    exit(0)
